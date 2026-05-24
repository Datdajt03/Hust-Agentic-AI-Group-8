import os, io, json, base64
from dotenv import load_dotenv
from PIL import Image
import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document
import zipfile
from google import genai
from google.genai import types
import mimetypes

load_dotenv()

client = genai.Client()

def call_vision_caption(image: Image.Image, prompt: str = "Describe the image content."):
    """
    Trả về caption (text) do Gemini 2.5 Flash tạo ra cho ảnh PIL.Image.
    """
    if image.mode in ("RGBA", "LA", "P"):
        image = image.convert("RGB")

    buffered = io.BytesIO()
    image.save(buffered, format="JPEG")
    image_bytes = buffered.getvalue()

    image_part = types.Part.from_bytes(data=image_bytes, mime_type="image/jpeg")

    response = client.models.generate_content(
        model="gemini-2.5-flash",
        contents=[image_part, prompt]
    )
    return response.text.strip()


def process_pdf(path):
    doc = fitz.open(path)
    parts = []
    for page in doc:
        text = page.get_text().strip()
        if text:
            parts.append(text)

        blocks = page.get_text("dict")["blocks"]
        for b in blocks:
            if b["type"] == 1 and "image" in b:
                img = Image.open(io.BytesIO(b["image"]))
                caption = call_vision_caption(
                    img,
                    prompt="Analyze the chart/diagram and list key observations."
                )
                parts.append(caption)
    return parts


def process_pptx(path):
    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
            if shape.shape_type == 13 and shape.image:
                img = Image.open(io.BytesIO(shape.image.blob))
                caption = call_vision_caption(
                    img,
                    prompt="Describe the diagram/chart and key points."
                )
                parts.append(caption)
    return parts


def process_docx(path):
    parts = []
    doc = Document(path)
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)

    with zipfile.ZipFile(path) as z:
        for f in z.namelist():
            if f.startswith("word/media/"):
                data = z.read(f)
                img = Image.open(io.BytesIO(data))
                caption = call_vision_caption(img, prompt="Describe this image content.")
                parts.append(caption)
    return parts


def process_txt(path):
    with open(path, "r", encoding="utf-8") as f:
        text = f.read().strip()
        return [text] if text else []


def process_file(path):
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        return process_pdf(path)
    if ext == ".pptx":
        return process_pptx(path)
    if ext == ".docx":
        return process_docx(path)
    if ext == ".txt":
        return process_txt(path)
    raise ValueError(f"Unsupported extension: {ext}")


def extract_content_as_string(input_path: str) -> str:
    """
    Xử lý file (PDF, PPTX, DOCX, TXT) và trả về một string duy nhất,
    gồm text và caption ghép lại.
    """
    parts = process_file(input_path)
    text_and_caption = "\n\n".join(parts)
    return text_and_caption
